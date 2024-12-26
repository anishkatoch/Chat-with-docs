import streamlit as st
import faiss
import os
from io import BytesIO
from docx import Document
import numpy as np
from PyPDF2 import PdfReader
from langchain.chains import RetrievalQA
from langchain.memory import ConversationSummaryBufferMemory
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.docstore.in_memory import InMemoryDocstore
from langchain.chat_models import ChatOpenAI
from dotenv import load_dotenv
from pptx import Presentation
# Load environment variables
load_dotenv()

# Set API key for OpenAI
openai_api_key = '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'

if not openai_api_key:
    st.error("OpenAI API Key is not set. Please add it to your environment variables.")
    st.stop()

MAX_MESSAGES = 5  # Maximum number of messages to retain


def answer_question(vectorstore, query, memory):
    """
    Answers a question based on the vectorstore and the conversation memory.
    Args:
        vectorstore: A FAISS vector store instance for retrieving context.
        query: The user's question/query.
        memory: A memory object storing the chat history.
    Returns:
        A string response from the model.
    """
    try:
        # Initialize the language model
        llm = ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key, temperature=0.7)

        # Setup the retriever using the vector store
        retriever = vectorstore.as_retriever(search_kwargs={"k": 2})  # Retrieve top 2 documents

        # Build the RetrievalQA chain with memory integration
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            retriever=retriever,
            memory=memory,  # Provides conversation context
            chain_type="stuff",  # Recommended default chain type
            return_source_documents=False,  # Turn off raw documents return (optional)
        )

        # Run the chain and return the output
        response = qa_chain.run(query)
        return response

    except ValueError as ve:
        raise ValueError(f"ValueError in answer_question: {ve}")
    except AttributeError as ae:
        raise AttributeError(f"AttributeError in answer_question: {ae}")
    except Exception as e:
        raise RuntimeError(f"Unexpected error in answer_question: {e}")



def process_input(input_data):
    """Processes different input types and returns a vectorstore."""
    texts = []
    for data in input_data:
        input_type, content = data

        if input_type == "PDF":
            pdf_reader = PdfReader(BytesIO(content.read()))
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
        elif input_type == "DOCX":
            doc = Document(BytesIO(content.read()))
            for para in doc.paragraphs:
                if para.text:
                    texts.append(para.text)
        elif input_type == "TXT":
            txt_content = content.read().decode("utf-8")
            texts.append(txt_content)
        elif input_type == "PPTX":
            ppt = Presentation(BytesIO(content.read()))
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)

    if not texts:
        raise ValueError("No text content extracted from the provided files.")

    # Use RecursiveCharacterTextSplitter for efficient splitting
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_texts = text_splitter.split_text(" ".join(texts))

    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)

    # Initialize FAISS vector store
    sample_embedding = np.array(embeddings.embed_query("sample text"))
    dimension = sample_embedding.shape[0]
    index = faiss.IndexFlatL2(dimension)

    # Create a vectorstore with batch processing
    vector_store = FAISS(
        embedding_function=embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={},
    )

    # Batch embed texts to avoid memory overload
    batch_size = 50
    for i in range(0, len(split_texts), batch_size):
        batch = split_texts[i:i + batch_size]
        vector_store.add_texts(batch)

    return vector_store

def truncate_messages(messages, max_messages=5):
    """
    Truncates the message list to retain only the last `max_messages` messages.
    If there are too many messages, summarizes the earlier ones.
    """
    if len(messages) > max_messages:
        # Summarize earlier messages (if needed)
        summary = " ".join(msg["content"] for msg in messages[:-max_messages])
        summarized_message = {
            "role": "assistant",
            "content": f"Summary of previous messages: {summary}"
        }
        # Keep only the summary and last few messages
        return [summarized_message] + messages[-max_messages:]
    return messages


# Modify the main function to handle PPTX uploads
def main():
    st.title("RAG Q&A App with GPT-4")

    with st.sidebar:
        st.header("Upload Inputs")

        pdf_files = st.file_uploader("Upload PDF Files", type="pdf", accept_multiple_files=True)
        docx_files = st.file_uploader("Upload DOCX Files", type="docx", accept_multiple_files=True)
        txt_files = st.file_uploader("Upload TXT Files", type="txt", accept_multiple_files=True)
        pptx_files = st.file_uploader("Upload PPTX Files", type="pptx", accept_multiple_files=True)

        if st.button("Process Input"):
            try:
                input_data = []
                if pdf_files:
                    input_data.extend([("PDF", file) for file in pdf_files])
                if docx_files:
                    input_data.extend([("DOCX", file) for file in docx_files])
                if txt_files:
                    input_data.extend([("TXT", file) for file in txt_files])
                if pptx_files:
                    input_data.extend([("PPTX", file) for file in pptx_files])

                if not input_data:
                    st.warning("Please provide at least one input.")
                    return

                vectorstore = process_input(input_data)
                st.session_state["vectorstore"] = vectorstore
                st.session_state["memory"] = ConversationSummaryBufferMemory(
                    llm=ChatOpenAI(model="gpt-4", openai_api_key=openai_api_key),
                    memory_key="chat_history",
                    max_token_limit=500  # Adjust token limit for summarization
                )
                st.session_state["messages"] = [
                    {"role": "assistant", "content": "Hello! How can I assist you today?"}
                ]
                st.success("Inputs processed successfully!")
            except Exception as e:
                st.error(f"Error processing input: {e}")

    if "vectorstore" in st.session_state and "memory" in st.session_state:
        if "messages" not in st.session_state:
            st.session_state["messages"] = [
                {"role": "assistant", "content": "Hello! How can I assist you today?"}
            ]

        # Truncate messages to prevent exceeding token limits
        st.session_state["messages"] = truncate_messages(st.session_state["messages"], MAX_MESSAGES)

        for msg in st.session_state.messages:
            st.chat_message(msg["role"]).write(msg["content"])

        if prompt := st.chat_input():
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.chat_message("user").write(prompt)

            # Truncate messages after adding a new one
            st.session_state.messages = truncate_messages(st.session_state.messages, MAX_MESSAGES)

            try:
                vectorstore = st.session_state["vectorstore"]
                memory = st.session_state["memory"]
                response = answer_question(vectorstore, prompt, memory)
                st.session_state.messages.append({"role": "assistant", "content": response})
                st.chat_message("assistant").write(response)
            except Exception as e:
                if "tokens" in str(e):
                    st.error("The input or output tokens exceeded the limit. Try shortening your query or context.")
                else:
                    st.error(f"Error generating response: {e}")


if __name__ == "__main__":
    main()